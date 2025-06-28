import json
import csv
import zipfile
import tarfile
import logging
from pathlib import Path
from typing import Optional, Dict, Callable
import chardet

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# File type handlers
class FileProcessor:
    def __init__(self):
        self.handlers: Dict[str, Callable] = {
            # Text files
            '.txt': self._read_text_file,
            '.md': self._read_text_file,
            '.rst': self._read_text_file,
            '.log': self._read_text_file,
            '.rtf': self._read_text_file,
            
            # Document files
            '.docx': self._read_docx,
            '.pdf': self._read_pdf,
            '.odt': self._read_odt,
            
            # Data files
            '.json': self._read_json,
            '.csv': self._read_csv,
            '.tsv': self._read_tsv,
            '.xml': self._read_xml,
            '.html': self._read_html,
            '.htm': self._read_html,
            
            # Code files
            '.py': self._read_text_file,
            '.js': self._read_text_file,
            '.java': self._read_text_file,
            '.cpp': self._read_text_file,
            '.c': self._read_text_file,
            '.h': self._read_text_file,
            '.css': self._read_text_file,
            '.sql': self._read_text_file,
            '.sh': self._read_text_file,
            '.yml': self._read_text_file,
            '.yaml': self._read_text_file,
            
            # Archive files
            '.zip': self._read_archive,
            '.tar': self._read_archive,
            '.gz': self._read_archive,
            
            # Office files
            '.xlsx': self._read_excel,
            '.xls': self._read_excel,
            '.pptx': self._read_powerpoint,
            
            # Image files (OCR)
            '.jpg': self._read_image,
            '.jpeg': self._read_image,
            '.png': self._read_image,
            '.bmp': self._read_image,
            '.tiff': self._read_image,
            '.gif': self._read_image,
        }
        
        # Import optional dependencies
        self._setup_optional_imports()
    
    def _setup_optional_imports(self):
        """Setup optional imports with fallbacks"""
        try:
            import openpyxl
            self.has_openpyxl = True
        except ImportError:
            self.has_openpyxl = False
            logger.warning("openpyxl not available - Excel files won't be supported")
        
        try:
            from bs4 import BeautifulSoup
            self.has_beautifulsoup = True
            self.BeautifulSoup = BeautifulSoup
        except ImportError:
            self.has_beautifulsoup = False
            logger.warning("beautifulsoup4 not available - HTML parsing will be basic")
        
        try:
            from pptx import Presentation
            self.has_pptx = True
            self.Presentation = Presentation
        except ImportError:
            self.has_pptx = False
            logger.warning("python-pptx not available - PowerPoint files won't be supported")

    def _detect_encoding(self, file_path: str) -> str:
        """Detect file encoding using chardet"""
        try:
            with open(file_path, 'rb') as file:
                raw_data = file.read(10000)  # Read first 10KB
                result = chardet.detect(raw_data)
                return result.get('encoding', 'utf-8') or 'utf-8'
        except Exception:
            return 'utf-8'

    def _read_text_file(self, file_path: str) -> str:
        """Read text files with encoding detection"""
        encoding = self._detect_encoding(file_path)
        encodings_to_try = [encoding, 'utf-8', 'latin-1', 'cp1252']
        
        for enc in encodings_to_try:
            try:
                with open(file_path, 'r', encoding=enc) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
        
        # Fallback: read as binary and decode w/o errors
        with open(file_path, 'rb') as file:
            return file.read().decode('utf-8', errors='ignore')

    def _read_docx(self, file_path: str) -> str:
        """Read DOCX files"""
        return docx2txt.process(file_path)

    def _read_pdf(self, file_path: str) -> str:
        """Read PDF files with proper cleanup"""
        doc = None
        try:
            doc = fitz.open(file_path)
            text_parts = []
            for page in doc:
                text_parts.append(page.get_text())
            return "\n".join(text_parts)
        finally:
            if doc:
                doc.close()

    def _read_odt(self, file_path: str) -> str:
        """Read ODT files (LibreOffice)"""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                content = zip_file.read('content.xml')
                root = ET.fromstring(content)
                # Extract text from ODT XML structures...
                text_parts = []
                for elem in root.iter():
                    if elem.text:
                        text_parts.append(elem.text)
                return ' '.join(text_parts)
        except Exception as e:
            logger.warning(f"Could not read ODT file {file_path}: {e}")
            return f"[ODT FILE] {os.path.basename(file_path)}"

    def _read_json(self, file_path: str) -> str:
        """Read JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
                return json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            # Fallback to reading as text
            return self._read_text_file(file_path)

    def _read_csv(self, file_path: str) -> str:
        """Read CSV files"""
        try:
            with open(file_path, 'r', newline='', encoding=self._detect_encoding(file_path)) as file:
                reader = csv.reader(file)
                rows = []
                for i, row in enumerate(reader):
                    if i > 1000:  # Limit rows for large files
                        rows.append("... (truncated)")
                        break
                    rows.append(' | '.join(row))
                return '\n'.join(rows)
        except Exception:
            return self._read_text_file(file_path)

    def _read_tsv(self, file_path: str) -> str:
        """Read TSV files"""
        try:
            with open(file_path, 'r', newline='', encoding=self._detect_encoding(file_path)) as file:
                reader = csv.reader(file, delimiter='\t')
                rows = []
                for i, row in enumerate(reader):
                    if i > 1000:
                        rows.append("... (truncated)")
                        break
                    rows.append(' | '.join(row))
                return '\n'.join(rows)
        except Exception:
            return self._read_text_file(file_path)

    def _read_xml(self, file_path: str) -> str:
        """Read XML files"""
        if self.has_beautifulsoup:
            try:
                with open(file_path, 'r', encoding=self._detect_encoding(file_path)) as file:
                    soup = self.BeautifulSoup(file.read(), 'xml')
                    return soup.get_text(separator=' ', strip=True)
            except Exception:
                pass
        return self._read_text_file(file_path)

    def _read_html(self, file_path: str) -> str:
        """Read HTML files"""
        if self.has_beautifulsoup:
            try:
                with open(file_path, 'r', encoding=self._detect_encoding(file_path)) as file:
                    soup = self.BeautifulSoup(file.read(), 'html.parser')
                    return soup.get_text(separator=' ', strip=True)
            except Exception:
                pass
        return self._read_text_file(file_path)

    def _read_excel(self, file_path: str) -> str:
        """Read Excel files"""
        if not self.has_openpyxl:
            return f"[EXCEL FILE] {os.path.basename(file_path)} (openpyxl not available)"
        
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"Sheet: {sheet_name}")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
                        
            workbook.close()
            return '\n'.join(text_parts)
        except Exception as e:
            logger.warning(f"Could not read Excel file {file_path}: {e}")
            return f"[EXCEL FILE] {os.path.basename(file_path)}"

    def _read_powerpoint(self, file_path: str) -> str:
        """Read PowerPoint files"""
        if not self.has_pptx:
            return f"[POWERPOINT FILE] {os.path.basename(file_path)} (python-pptx not available)"
        
        try:
            prs = self.Presentation(file_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides):
                text_parts.append(f"Slide {i + 1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
                        
            return '\n'.join(text_parts)
        except Exception as e:
            logger.warning(f"Could not read PowerPoint file {file_path}: {e}")
            return f"[POWERPOINT FILE] {os.path.basename(file_path)}"

    def _read_archive(self, file_path: str) -> str:
        """Read archive files (list contents)"""
        try:
            if file_path.endswith('.zip'):
                with zipfile.ZipFile(file_path, 'r') as archive:
                    file_list = archive.namelist()
            elif file_path.endswith(('.tar', '.tar.gz', '.tgz')):
                with tarfile.open(file_path, 'r') as archive:
                    file_list = archive.getnames()
            else:
                return f"[ARCHIVE FILE] {os.path.basename(file_path)}"
            
            return f"[ARCHIVE FILE] {os.path.basename(file_path)}\nContents:\n" + '\n'.join(file_list[:100])
        except Exception as e:
            logger.warning(f"Could not read archive {file_path}: {e}")
            return f"[ARCHIVE FILE] {os.path.basename(file_path)}"

    def _read_image(self, file_path: str) -> str:
        """Read images using OCR"""
        try:
            return pytesseract.image_to_string(Image.open(file_path))
        except Exception as e:
            logger.warning(f"OCR failed for {file_path}: {e}")
            return f"[IMAGE FILE] {os.path.basename(file_path)}"




def open_file(file_path: str) -> Optional[str]:
    """
    File Reading, with a support for multiple mime types..
    Args:
        file_path (str): Path to the file to read
    Returns:
        str or None: Extracted text content or None if processing failed
    """
    file_processor = FileProcessor()
    # Validation checks
    path_obj = Path(file_path)
    if not path_obj.exists():
        logger.error(f"File not found: {file_path}")
        return None
    
    if not path_obj.is_file():
        logger.error(f"Path is not a file: {file_path}")
        return None
    
    # Size check (50MB limit)
    file_size = path_obj.stat().st_size
    max_size = 50 * 1024 * 1024  # 50MB
    if file_size > max_size:
        logger.error(f"File too large: {file_path} ({file_size:,} bytes > {max_size:,} bytes)")
        return None
    
    if file_size == 0:
        logger.warning(f"Empty file: {file_path}")
        return ""
    
    try:
        # Get File Exten...
        file_ext = path_obj.suffix.lower()
        
        # Check if we have a specific handler
        if file_ext in file_processor.handlers:
            logger.info(f"Processing {file_path} as {file_ext}")
            return file_processor.handlers[file_ext](file_path)
        
        # Fallback: use MIME type
        mime_type, _ = mimetypes.guess_type(file_path)
        
        if mime_type:
            if mime_type.startswith("text/"):
                return file_processor._read_text_file(file_path)
            elif mime_type.startswith("image/"):
                return file_processor._read_image(file_path)
            elif mime_type.startswith("audio/"):
                return f"[AUDIO FILE] {path_obj.name}"
            elif mime_type.startswith("video/"):
                return f"[VIDEO FILE] {path_obj.name}"
        
        # Final fallback: try as text
        logger.warning(f"Unknown file type for {file_path}, attempting to read as text")
        return file_processor._read_text_file(file_path)
        
    except Exception as e:
        logger.error(f"Failed to read file {file_path}: {e}")
        return None